from fastapi import FastAPI, UploadFile, Form, File, HTTPException
from typing import List, Optional
from pydantic import BaseModel
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct
import requests, uuid, os, base64
from io import BytesIO
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

# ---- Env config ----
OLLAMA = os.getenv("OLLAMA_BASE_URL", "http://ollama:11434")
EMBED_MODEL = os.getenv("EMBED_MODEL", "nomic-embed-text:latest")
GEN_MODEL = os.getenv("GEN_MODEL", "llama3.1:8b")
QURL = os.getenv("QDRANT_URL", "http://localhost:6333")
COLL = os.getenv("COLLECTION", "docs")
CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "800"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "120"))

# runtime-tunable
EMBED_BATCH = int(os.getenv("EMBED_BATCH", "32"))
UPSERT_BATCH = int(os.getenv("UPSERT_BATCH", "256"))
RUNTIME = {
    "EMBED_BATCH": EMBED_BATCH,
    "UPSERT_BATCH": UPSERT_BATCH,
    "CHUNK_SIZE": CHUNK_SIZE,
    "CHUNK_OVERLAP": CHUNK_OVERLAP,
}

app = FastAPI()
qdrant = QdrantClient(url=QURL)

def ensure_collection(dim=768):
    names = [c.name for c in qdrant.get_collections().collections]
    if COLL not in names:
        qdrant.recreate_collection(
            collection_name=COLL,
            vectors_config=VectorParams(size=dim, distance=Distance.COSINE),
        )

def embed_batch(text_list):
    if not text_list:
        return []
    embeddings = []
    for text in text_list:
        if text and text.strip():  # Skip empty strings
            r = requests.post(f"{OLLAMA}/api/embeddings", json={
                "model": EMBED_MODEL,
                "prompt": text
            }, timeout=300)
            r.raise_for_status()
            data = r.json()
            embeddings.append(data["embedding"])
        else:
            # For empty strings, create a zero vector (768 dimensions for nomic-embed-text)
            embeddings.append([0.0] * 768)
    return embeddings

def embed(text):
    r = requests.post(f"{OLLAMA}/api/embeddings", json={
        "model": EMBED_MODEL, "prompt": text
    }, timeout=120)
    r.raise_for_status()
    data = r.json()
    return data["embedding"]

def chunk_list(items, n):
    for i in range(0, len(items), n):
        yield items[i:i+n]

def chunks(s: str, size: int, overlap: int):
    out, i = [], 0
    if not s: return out
    step = max(1, size - overlap)
    while i < len(s):
        out.append(s[i:i+size])
        i += step
    return out

# ---------- File extractors ----------
def extract_text_from_txt(b: bytes) -> str:
    return b.decode("utf-8", errors="ignore")

def extract_text_from_pdf(b: bytes) -> str:
    bio = BytesIO(b)
    reader = PdfReader(bio)
    pages = []
    for p in reader.pages:
        try:
            pages.append(p.extract_text() or "")
        except Exception:
            pages.append("")
    return "\n\n".join(pages).strip()

def extract_text_from_docx(b: bytes) -> str:
    bio = BytesIO(b)
    doc = DocxDocument(bio)
    lines = []
    for para in doc.paragraphs:
        if para.text:
            lines.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines).strip()

def extract_text_from_pptx(b: bytes) -> str:
    bio = BytesIO(b)
    prs = Presentation(bio)
    slides_text = []
    for i, slide in enumerate(prs.slides, start=1):
        tparts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                tparts.append(shape.text)
        if tparts:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(tparts))
    return "\n\n".join(slides_text).strip()

def extract_text_from_csv(b: bytes) -> str:
    bio = BytesIO(b)
    df = pd.read_csv(bio)
    return df.to_csv(index=False)

def extract_text_from_xlsx(b: bytes) -> str:
    bio = BytesIO(b)
    xls = pd.read_excel(bio, sheet_name=None, engine="openpyxl")
    parts = []
    for name, df in xls.items():
        parts.append(f"[Sheet: {name}]")
        parts.append(df.to_csv(index=False))
    return "\n\n".join(parts).strip()

def extract_text_from_xls(b: bytes) -> str:
    bio = BytesIO(b)
    xls = pd.read_excel(bio, sheet_name=None, engine="xlrd")
    parts = []
    for name, df in xls.items():
        parts.append(f"[Sheet: {name}]")
        parts.append(df.to_csv(index=False))
    return "\n\n".join(parts).strip()

def extract_text(filename: str, content_type: str, b: bytes) -> str:
    name = (filename or "").lower()
    ct = (content_type or "").lower()
    try:
        if name.endswith(".pdf") or "pdf" in ct:
            return extract_text_from_pdf(b)
        if name.endswith(".docx") or "word" in ct:
            return extract_text_from_docx(b)
        if name.endswith(".pptx") or "powerpoint" in ct:
            return extract_text_from_pptx(b)
        if name.endswith(".csv") or "csv" in ct:
            return extract_text_from_csv(b)
        if name.endswith(".xlsx"):
            return extract_text_from_xlsx(b)
        if name.endswith(".xls"):
            return extract_text_from_xls(b)
        if name.endswith(".md") or name.endswith(".txt") or "text" in ct:
            return extract_text_from_txt(b)
    except Exception:
        try:
            return b.decode("utf-8", errors="ignore")
        except Exception:
            return ""
    return b.decode("utf-8", errors="ignore")

# ---------- API ----------
@app.on_event("startup")
def _start():
    v = embed("dimension probe")
    ensure_collection(len(v))

@app.get("/config")
def get_config():
    return {**RUNTIME, "EMBED_MODEL": EMBED_MODEL, "GEN_MODEL": GEN_MODEL, "COLLECTION": COLL}

class ConfigPatch(BaseModel):
    EMBED_BATCH: Optional[int] = None
    UPSERT_BATCH: Optional[int] = None
    CHUNK_SIZE: Optional[int] = None
    CHUNK_OVERLAP: Optional[int] = None

@app.post("/config")
def set_config(p: ConfigPatch):
    if p.EMBED_BATCH is not None and p.EMBED_BATCH > 0:
        RUNTIME["EMBED_BATCH"] = int(p.EMBED_BATCH)
    if p.UPSERT_BATCH is not None and p.UPSERT_BATCH > 0:
        RUNTIME["UPSERT_BATCH"] = int(p.UPSERT_BATCH)
    if p.CHUNK_SIZE is not None and p.CHUNK_SIZE >= 128:
        RUNTIME["CHUNK_SIZE"] = int(p.CHUNK_SIZE)
    if p.CHUNK_OVERLAP is not None and 0 <= p.CHUNK_OVERLAP < RUNTIME["CHUNK_SIZE"]:
        RUNTIME["CHUNK_OVERLAP"] = int(p.CHUNK_OVERLAP)
    return {"ok": True, **RUNTIME}

def http_get(url: str, bearer: Optional[str] = None) -> bytes:
    headers = {}
    if bearer:
        headers["Authorization"] = f"Bearer {bearer}"
    r = requests.get(url, headers=headers, timeout=300)
    if r.status_code != 200:
        raise HTTPException(status_code=502, detail=f"Fetch failed {r.status_code} for {url}")
    return r.content

def ingest_bytes(filename: str, content_type: str, raw: bytes, project: str, chunks_out: list, meta_out: list):
    text = extract_text(filename, content_type, raw)
    for ch in chunks(text, RUNTIME["CHUNK_SIZE"], RUNTIME["CHUNK_OVERLAP"]):
        chunks_out.append(ch)
        meta_out.append({"project": project, "file": filename})

@app.post("/ingest")
async def ingest(files: List[UploadFile] = File(...), project: str = Form("default")):
    all_chunks, meta = [], []
    for f in files:
        b = await f.read()
        ingest_bytes(f.filename, f.content_type or "", b, project, all_chunks, meta)

    total_points = 0
    for ch_batch in chunk_list(all_chunks, RUNTIME["EMBED_BATCH"]):
        vecs = embed_batch(ch_batch)
        pts = []
        base_idx = total_points  # not used for mapping, but okay
        for i, v in enumerate(vecs):
            m = meta[(total_points + i) if (total_points + i) < len(meta) else -1]
            pts.append(PointStruct(id=str(uuid.uuid4()), vector=v, payload={"text": ch_batch[i], **m}))
        for up in chunk_list(pts, RUNTIME["UPSERT_BATCH"]):
            qdrant.upsert(collection_name=COLL, points=up)
            total_points += len(up)

    return {"ok": True, "project": project, "chunks": total_points, "files": [f.filename for f in files]}

class UrlIngest(BaseModel):
    project: str = "default"
    urls: List[str]
    bearer: Optional[str] = None

@app.post("/ingest_urls")
def ingest_urls(body: UrlIngest):
    all_chunks, meta = [], []
    fetched = 0
    for u in body.urls:
        try:
            data = http_get(u, body.bearer)
            lower = u.lower()
            ctype = "application/octet-stream"
            if lower.endswith(".pdf"): ctype = "application/pdf"
            elif lower.endswith(".docx"): ctype = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            elif lower.endswith(".pptx"): ctype = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            elif lower.endswith(".xlsx"): ctype = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            elif lower.endswith(".xls"): ctype = "application/vnd.ms-excel"
            elif lower.endswith(".csv"): ctype = "text/csv"
            elif lower.endswith(".md") or lower.endswith(".txt"): ctype = "text/plain"
            fname = u.split("?")[0].rstrip("/").split("/")[-1] or "downloaded"
            ingest_bytes(fname, ctype, data, body.project, all_chunks, meta)
            fetched += 1
        except Exception as e:
            print("URL ingest error:", u, e)
            continue

    total_points = 0
    for ch_batch in chunk_list(all_chunks, RUNTIME["EMBED_BATCH"]):
        vecs = embed_batch(ch_batch)
        pts = [PointStruct(id=str(uuid.uuid4()), vector=v,
                           payload={"text": ch_batch[i], **meta[total_points + i]}) for i, v in enumerate(vecs)]
        for up in chunk_list(pts, RUNTIME["UPSERT_BATCH"]):
            qdrant.upsert(collection_name=COLL, points=up)
            total_points += len(up)

    return {"ok": True, "project": body.project, "fetched": fetched, "chunks": total_points}

class ShareIngest(BaseModel):
    project: str = "default"
    share_url: str
    bearer: str

GRAPH_BASE = os.getenv("GRAPH_BASE", "https://graph.microsoft.com/v1.0")

def graph_headers(token: str):
    return {"Authorization": f"Bearer {token}"}

def encode_share_url(u: str) -> str:
    b64 = base64.b64encode(u.encode()).decode()
    return "u!" + b64.rstrip("=")

@app.post("/ingest_sharepoint")
def ingest_sharepoint(b: ShareIngest):
    sid = encode_share_url(b.share_url)
    r = requests.get(f"{GRAPH_BASE}/shares/{sid}/driveItem", headers=graph_headers(b.bearer), timeout=60)
    if r.status_code != 200:
        raise HTTPException(status_code=502, detail=f"Graph resolve failed: {r.text}")
    item = r.json()

    if "folder" in item:
        rid = item["id"]
        parent_drive = item["parentReference"]["driveId"]
        r2 = requests.get(f"{GRAPH_BASE}/drives/{parent_drive}/items/{rid}/children",
                          headers=graph_headers(b.bearer), timeout=120)
        if r2.status_code != 200:
            raise HTTPException(status_code=502, detail=f"Graph list failed: {r2.text}")
        children = [c for c in r2.json().get("value", []) if "file" in c]
    else:
        children = [item] if "file" in item else []

    all_chunks, meta = [], []
    fetched = 0
    for c in children:
        drive_id = c["parentReference"]["driveId"]
        item_id  = c["id"]
        name     = c.get("name","file")
        r3 = requests.get(f"{GRAPH_BASE}/drives/{drive_id}/items/{item_id}/content",
                          headers=graph_headers(b.bearer), timeout=300)
        if r3.status_code != 200:
            print("Graph download failed:", name, r3.text); continue
        raw = r3.content
        lower = name.lower()
        ctype = "application/octet-stream"
        if lower.endswith(".pdf"): ctype = "application/pdf"
        elif lower.endswith(".docx"): ctype = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif lower.endswith(".pptx"): ctype = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif lower.endswith(".xlsx"): ctype = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        elif lower.endswith(".xls"): ctype = "application/vnd.ms-excel"
        elif lower.endswith(".csv"): ctype = "text/csv"
        elif lower.endswith(".md") or lower.endswith(".txt"): ctype = "text/plain"
        ingest_bytes(name, ctype, raw, b.project, all_chunks, meta)
        fetched += 1

    total_points = 0
    for ch_batch in chunk_list(all_chunks, RUNTIME["EMBED_BATCH"]):
        vecs = embed_batch(ch_batch)
        pts = [PointStruct(id=str(uuid.uuid4()), vector=v,
                           payload={"text": ch_batch[i], **meta[total_points + i]}) for i, v in enumerate(vecs)]
        for up in chunk_list(pts, RUNTIME["UPSERT_BATCH"]):
            qdrant.upsert(collection_name=COLL, points=up)
            total_points += len(up)

    return {"ok": True, "project": b.project, "fetched": fetched, "chunks": total_points}

class Query(BaseModel):
    question: str
    top_k: int = 4
    project: str = "default"

@app.post("/chat")
def chat(q: Query):
    vec = embed(q.question)
    # filter by project in payload when qdrant supports filter here (simplified: no filter)
    res = qdrant.search(collection_name=COLL, query_vector=vec, limit=q.top_k, with_payload=True)
    ctx = "\n\n".join([r.payload.get("text","") for r in res])
    prompt = (
        "You are a helpful assistant. Use the context to answer concisely. "
        "If the answer is not in the context, say you don't know.\n\n"
        f"Context:\n{ctx}\n\nQuestion: {q.question}\nAnswer:"
    )
    r = requests.post(f"{OLLAMA}/api/generate", json={
        "model": GEN_MODEL, "prompt": prompt, "stream": False, "options": {"num_ctx": 1024}
    }, timeout=300)
    r.raise_for_status()
    ans = r.json().get("response", "")
    return {"answer": ans, "sources": [{"score": float(hit.score), "file": hit.payload.get("file")} for hit in res]}
